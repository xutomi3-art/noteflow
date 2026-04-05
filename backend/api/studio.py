import io
import json as json_lib
import logging
import os
import re
import tempfile
import time
import urllib.parse
import uuid

from pydub import AudioSegment

from fastapi import APIRouter, Body, Depends, HTTPException
from fastapi.responses import StreamingResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from backend.core.config import settings
from backend.core.database import get_db
from backend.core.deps import get_current_user
from backend.models.notebook import Notebook
from backend.models.source import Source
from backend.models.user import User
from pydantic import BaseModel

from backend.services.qwen_client import qwen_client
from backend.services.ragflow_client import ragflow_client
from backend.services import permission_service
from backend.services.tts_client import text_to_speech
from backend.services.docmee_client import docmee_client

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/notebooks/{notebook_id}/studio", tags=["studio"])
ppt_router = APIRouter(prefix="/ppt", tags=["ppt"])

def _detect_language(text: str) -> str:
    """Detect if text is primarily Chinese or English based on character ratio."""
    cjk_count = sum(1 for ch in text if '\u4e00' <= ch <= '\u9fff')
    return "Chinese" if cjk_count > len(text) * 0.05 else "English"


PROMPTS = {
    "summary": """Based on the following document contents, write a comprehensive summary that covers all key topics and main points. Structure it with clear sections and bullet points. 
Formatting rules:
- Use ## for main section headers and ### for sub-sections. Do NOT use #### or deeper headings.
- Use bullet points (- ) for details under each section.
- Keep the structure flat and readable.

DOCUMENTS:
{context}""",
    "faq": """Based on the following document contents, generate a list of 8-10 frequently asked questions with detailed answers. Each Q&A should cover an important concept from the documents. Write in the same language as the documents. Format as:

Q: [question]
A: [answer]

DOCUMENTS:
{context}""",
    "study_guide": """Based on the following document contents, create a comprehensive study guide that includes:
1. Key concepts and definitions
2. Important relationships between ideas
3. Summary of each major section
4. Review questions for self-assessment


DOCUMENTS:
{context}""",
    "action_items": """Based on the following document contents, extract all action items, tasks, to-dos, next steps, and follow-up items. For each action item, include:
- The specific task or action required
- Who is responsible (if mentioned)
- Deadline or timeline (if mentioned)
- Priority level (High/Medium/Low) based on context

Group related action items together under clear category headers. 
DOCUMENTS:
{context}""",
    "mindmap": """Generate a mind map JSON structure from the source documents. Node labels should be in the same language as the documents.
Return ONLY valid JSON in this exact format:
{{
  "nodes": [
    {{"id": "root", "label": "Central Topic", "level": 0}},
    {{"id": "n1", "label": "Main Branch 1", "level": 1, "parent": "root"}},
    {{"id": "n1_1", "label": "Sub Topic", "level": 2, "parent": "n1"}}
  ]
}}

Rules:
1. Root node has level=0 and no parent field
2. Main branches have level=1 and parent="root"
3. Sub-topics have level=2 and parent=the level-1 node id
4. Max 5 level-1 nodes, max 3 level-2 nodes per level-1 node
5. Labels are concise (under 6 words)
6. Each node has a unique id string
7. Return ONLY valid JSON, no markdown fences or explanation

Source content:
{context}""",
    "swot": """You are a strategic analyst. Based on the following documents (which may include meeting transcripts, reports, or other materials), perform a thorough SWOT analysis on the core topic discussed.

Output format:

## 📊 SWOT Analysis: [auto-extract core topic]

**Context**: (2-3 sentences summarizing what is being analyzed)

### 💪 Strengths
- **[Title]**: Specific explanation with evidence from the documents
- (3-5 items)

### ⚠️ Weaknesses
- **[Title]**: Specific explanation with evidence from the documents
- (3-5 items)

### 🚀 Opportunities
- **[Title]**: Specific explanation referencing trends or possibilities mentioned
- (3-5 items)

### 🔴 Threats
- **[Title]**: Specific explanation with evidence from the documents
- (3-5 items)

### 📌 Recommendations
Based on the SWOT analysis, provide 3-5 strategic recommendations, each with:
- The recommendation
- Which SWOT factor it addresses
- Priority (High/Medium/Low)

Rules:
- Every item must have specific evidence from the documents, no generic statements
- Preserve all specific data, names, amounts, dates
- If information is insufficient, state clearly what is missing
- Write in the same language as the documents

DOCUMENTS:
{context}""",
    "recommendations": """You are an experienced business consultant. Based on the following documents (which may include meeting transcripts, reports, or other materials), provide detailed, actionable recommendations for the issues and challenges discussed.

Output format:

## 💡 Recommendations

**Issue Overview**: (3-5 sentences summarizing the problems and challenges)

### Recommendation 1: [Title]
- **Problem**: What specific issue this addresses (cite from documents)
- **Solution**: Step-by-step action plan (be as specific as possible)
- **Rationale**: Why this approach (backed by document data or domain knowledge)
- **Expected Outcome**: What results to expect
- **Risks**: Potential obstacles

### Recommendation 2: [Title]
(same format)

### Recommendation 3: [Title]
(same format)

(Provide 3-5 core recommendations)

### Quick Action Items
| # | Action | Suggested Owner | Timeline | Priority |
|---|--------|----------------|----------|----------|
| 1 | ... | ... | ... | High/Med/Low |

### Issues Requiring Further Discussion
- Points where **no consensus** was reached, with each party's position
- Topics **not covered but should be discussed**

Rules:
- Recommendations must be specific and actionable, not vague like "improve communication"
- Each recommendation must be supported by evidence from the documents
- Preserve all specific data, names, amounts, dates
- Write in the same language as the documents

DOCUMENTS:
{context}""",
    "risk_analysis": """You are a professional risk management consultant. Based on the following documents (which may include meeting transcripts, reports, or other materials), identify and analyze all risks discussed or implied.

Output format:

## 🔍 Risk Analysis Report

**Scope**: (2-3 sentences describing what is being analyzed)

### Identified Risks

**🔴 High Risk**

**Risk 1: [Name]**
- **Description**: What the risk is (cite from documents)
- **Trigger**: Under what conditions it would occur
- **Impact**: Which areas affected (financial/timeline/personnel/compliance)
- **Severity**: Specific potential consequences
- **Current Status**: Any existing mitigation discussed
- **Recommended Action**: Specific risk mitigation steps

**🟡 Medium Risk**

**Risk 2: [Name]**
(same format, may be slightly abbreviated)

**🟢 Low Risk**

**Risk 3: [Name]**
(brief description)

### Potential Risks (not discussed but noteworthy)
- Risks **not mentioned in the documents but likely relevant** based on context
- Explain why each is considered a potential risk

### Risk Matrix
| Risk | Likelihood | Impact | Level | Urgency |
|------|-----------|--------|-------|---------|
| ... | High/Med/Low | High/Med/Low | 🔴/🟡/🟢 | Immediate/Short-term/Long-term |

### Recommended Next Steps
- Risks requiring **immediate action** with specific steps
- Risks requiring **ongoing monitoring** with suggested indicators
- Risks requiring **further assessment** with evaluation methods

Rules:
- Risks must be specific, not generic like "market risk"
- Each risk must have evidence from the documents
- Distinguish between "known risks" and "potential risks"
- Preserve all specific data, names, amounts, dates
- Write in the same language as the documents

DOCUMENTS:
{context}""",
    "decision_support": """You are a decision analysis expert. Based on the following documents (which may include meeting transcripts, reports, or other materials), identify all items requiring decisions and provide structured decision support analysis.

Output format:

## 🎯 Decision Support Analysis

**Overview**: (Summarize which items need decisions, which are decided, which are pending)

### Decisions Already Made
| # | Decision | Key Rationale | Owner | Notes |
|---|---------|--------------|-------|-------|
| 1 | ... | ... | ... | ... |

### Pending Decision Analysis

**Pending Item 1: [Decision Question]**

**Background**: Why this decision is needed (cite from documents)

**Options Comparison**:
| Dimension | Option A: [Name] | Option B: [Name] | Option C: [Name] |
|-----------|-----------------|-----------------|-----------------|
| Approach | ... | ... | ... |
| Pros | ... | ... | ... |
| Cons | ... | ... | ... |
| Cost | ... | ... | ... |
| Timeline | ... | ... | ... |
| Risk | ... | ... | ... |

**Stakeholder Positions**:
- [Person/Role A] leans toward: ... Reason: ...
- [Person/Role B] leans toward: ... Reason: ...

**AI Recommendation**:
- Recommended option: [X]
- Rationale: (backed by document data and discussion content)
- Watch out for: ...

**Pending Item 2: [Decision Question]**
(same format)

### Implicit Decision Points
- Items **implied but not explicitly raised** that need decisions
- Why each also requires a decision

### Decision Timeline
| Pending Item | Suggested Deadline | Reason | Info Needed Before Deciding |
|-------------|-------------------|--------|---------------------------|
| ... | ... | ... | ... |

Rules:
- Options comparison must be fair and objective, not biased
- AI recommendation must be clear, not just "each has pros and cons"
- If insufficient info for a recommendation, state what info is missing
- Quote original statements as evidence using 「」
- Preserve all specific data, names, amounts, dates
- Write in the same language as the documents

DOCUMENTS:
{context}""",
}

PPT_PROMPT = """You are creating a professional PowerPoint presentation based on the source documents.
Generate a JSON structure for a presentation.

Configuration:
- Number of slides: {n_slides}
- Tone: {tone}
- Verbosity: {verbosity}
- Language: {language}

JSON format:
{{
  "title": "Presentation Title",
  "subtitle": "A brief one-line subtitle",
  "slides": [
    {{
      "title": "Slide Title",
      "bullets": [
        {{"main": "Key point", "sub": "Supporting detail or explanation"}},
        {{"main": "Key point 2", "sub": "More context here"}}
      ]
    }}
  ]
}}

Rules:
1. Generate exactly {n_slides} content slides (not counting the title slide)
2. IMPORTANT: Vary the number of bullet points per slide to look natural and human-made. Some slides should have 2 bullets, some 3, some 4, and occasionally 5. Do NOT give every slide the same number of bullets.
3. "main" is a bold headline (under 10 words), "sub" is a supporting sentence
4. If verbosity is "concise", keep "sub" very short (under 8 words) or empty string
5. If verbosity is "text-heavy", make "sub" a full explanatory sentence
6. Match the tone: casual=conversational, professional=formal, funny=witty, educational=instructive
7. Write in {language}
8. Make the presentation feel like a real human created it — vary structure, emphasis, and depth across slides
9. Return ONLY valid JSON, no markdown fences

Source content:
{context}"""

PODCAST_PROMPT = """Create a natural podcast dialogue between Host and Guest based on the source documents.
Format each line strictly as:
HOST: [what the host says]
GUEST: [what the guest says]

Rules:
1. Alternate between HOST and GUEST, starting with HOST
2. 8-12 exchanges total
3. HOST introduces the topic, GUEST provides insights
4. Keep each line under 50 words
5. Make it conversational and engaging
6. Return ONLY the dialogue, no titles or descriptions

Source content:
{context}"""

# Broad retrieval queries to get comprehensive coverage of all document content
# Bilingual to handle both English and Chinese documents
_RETRIEVAL_QUERIES = [
    "main topics key points overview 主要内容 关键要点 概述",
    "important details information content 重要信息 详细内容",
    "background context introduction 背景介绍 基本概念",
]


async def _get_source_context(db: AsyncSession, notebook_id: uuid.UUID, source_ids: list[str] | None = None) -> str:
    """Get document content for studio generation.

    Strategy:
    1. For TXT/MD: read directly from local storage (fast, no RAGFlow round-trip)
    2. For PDF/DOCX/PPTX: retrieve chunks from RAGFlow (MinerU-parsed content lives there)
    3. Combine both with deduplication

    If source_ids is provided, only include those sources. Otherwise use all ready sources.
    """
    query = select(Source).where(
        Source.notebook_id == notebook_id,
        Source.status == "ready",
    )
    if source_ids:
        query = query.where(Source.id.in_([uuid.UUID(sid) for sid in source_ids]))
    result = await db.execute(query)
    sources = list(result.scalars().all())
    if not sources:
        return ""

    context_parts: list[str] = []
    has_ragflow_sources = False

    # Step 1: Read TXT/MD files directly
    for source in sources:
        if source.file_type in ("txt", "md") and source.storage_url:
            try:
                with open(source.storage_url, "r", encoding="utf-8", errors="replace") as f:
                    content = f.read()[:8000]
                if content.strip():
                    context_parts.append(f"--- {source.filename} ---\n{content}")
            except Exception:
                pass
        elif source.ragflow_doc_id:
            has_ragflow_sources = True

    # Step 2: Retrieve chunks from RAGFlow for PDF/DOCX/PPTX sources
    if has_ragflow_sources:
        nb_result = await db.execute(
            select(Notebook).where(Notebook.id == notebook_id)
        )
        notebook = nb_result.scalar_one_or_none()
        dataset_id = notebook.ragflow_dataset_id if notebook else None

        # Filter RAGFlow retrieval to selected sources only
        filter_doc_ids = [s.ragflow_doc_id for s in sources if s.ragflow_doc_id] if source_ids else None

        if dataset_id:
            seen_texts: set[str] = set()
            all_chunks: list[dict] = []

            # Use multiple broad queries to get comprehensive coverage
            for query in _RETRIEVAL_QUERIES:
                try:
                    chunks = await ragflow_client.retrieve(
                        dataset_ids=[dataset_id],
                        question=query,
                        top_k=settings.RAG_TOP_K,
                        document_ids=filter_doc_ids,
                    )
                except Exception:
                    chunks = []
                for chunk in chunks:
                    text = chunk.get("content_with_weight", chunk.get("content", "")).strip()
                    if text and text not in seen_texts:
                        seen_texts.add(text)
                        all_chunks.append(chunk)

            if all_chunks:
                # Group chunks by source document
                doc_chunks: dict[str, list[str]] = {}
                for chunk in all_chunks:
                    doc_name = chunk.get("document_keyword", chunk.get("docnm_kwd", "document"))
                    text = chunk.get("content_with_weight", chunk.get("content", "")).strip()
                    if text:
                        doc_chunks.setdefault(doc_name, []).append(text)

                for doc_name, texts in doc_chunks.items():
                    combined = "\n\n".join(texts)[:10000]
                    context_parts.append(f"--- {doc_name} ---\n{combined}")

    return "\n\n".join(context_parts)


class PptGenerateRequest(BaseModel):
    template_id: str = ""
    scene: str = ""
    audience: str = ""
    language: str = "zh"
    length: str = "medium"


_FALLBACK_GENERATION_OPTIONS = {
    "lang": [
        {"label": "English", "value": "en"},
        {"label": "Chinese", "value": "zh"},
        {"label": "Japanese", "value": "ja"},
        {"label": "Korean", "value": "ko"},
    ],
    "scene": [
        {"label": "Work Report", "value": "work_report"},
        {"label": "Education", "value": "education"},
        {"label": "Business Plan", "value": "business_plan"},
        {"label": "Product Introduction", "value": "product_intro"},
        {"label": "Academic Research", "value": "academic"},
        {"label": "Project Summary", "value": "project_summary"},
    ],
    "audience": [
        {"label": "Team / Colleagues", "value": "team"},
        {"label": "Management", "value": "management"},
        {"label": "Clients", "value": "client"},
        {"label": "Students", "value": "student"},
        {"label": "General", "value": "general"},
    ],
}


@ppt_router.get("/templates")
async def list_ppt_templates(
    page: int = 1,
    size: int = 20,
    lang: str = "",
    current_user: User = Depends(get_current_user),
):
    """List available Docmee PPT templates."""
    result = await docmee_client.list_templates(page=page, size=size, lang=lang)
    return result


@ppt_router.get("/template-options")
async def get_ppt_template_options(
    current_user: User = Depends(get_current_user),
):
    """Get template filter options (categories, etc.)."""
    return await docmee_client.get_template_options()


@ppt_router.get("/generation-options")
async def get_ppt_generation_options(
    current_user: User = Depends(get_current_user),
):
    """Get generation options (scene, audience, language). Falls back to defaults when Docmee is unavailable."""
    result = await docmee_client.get_generation_options()
    if not result or (not result.get("lang") and not result.get("scene") and not result.get("audience")):
        return _FALLBACK_GENERATION_OPTIONS

    # Docmee returns 'name' instead of 'label' — normalize for frontend.
    # Use Docmee lang options (native names are fine) but English fallbacks for scene/audience.
    normalized: dict = {}
    lang_items = result.get("lang", [])
    if lang_items:
        normalized["lang"] = [{"label": item.get("name", item.get("label", "")), "value": item.get("value", "")} for item in lang_items]
    else:
        normalized["lang"] = _FALLBACK_GENERATION_OPTIONS["lang"]
    normalized["scene"] = _FALLBACK_GENERATION_OPTIONS["scene"]
    normalized["audience"] = _FALLBACK_GENERATION_OPTIONS["audience"]
    return normalized


@router.post("/ppt")
async def generate_ppt(
    notebook_id: uuid.UUID,
    config: PptGenerateRequest | None = None,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    if not await permission_service.check_permission(db, notebook_id, current_user.id, "view"):
        raise HTTPException(status_code=403, detail="No access to this notebook")
    context = await _get_source_context(db, notebook_id)
    if not context:
        raise HTTPException(status_code=400, detail="No source documents available")

    # Get notebook name for the download filename
    nb_result = await db.execute(select(Notebook).where(Notebook.id == notebook_id))
    notebook = nb_result.scalar_one_or_none()
    nb_name = notebook.name if notebook else "presentation"
    safe_name = re.sub(r'[^\w\s-]', '', nb_name).strip()
    safe_name = re.sub(r'[\s]+', '_', safe_name) or 'presentation'
    filename = f"{safe_name}.pptx"

    cfg = config or PptGenerateRequest()

    # Try Docmee first
    if await docmee_client.is_available() and cfg.template_id:
        ppt_info = await docmee_client.generate_ppt(
            content=context[:1000],
            template_id=cfg.template_id,
            scene=cfg.scene,
            audience=cfg.audience,
            lang=cfg.language,
            length=cfg.length,
        )
        if ppt_info:
            ppt_id = ppt_info.get("id", "")
            pptx_bytes = await docmee_client.download_pptx(ppt_id, lang=cfg.language)
            if pptx_bytes:
                buf = io.BytesIO(pptx_bytes)
                return StreamingResponse(
                    buf,
                    media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    headers={"Content-Disposition": f'attachment; filename="{urllib.parse.quote(filename)}"; filename*=UTF-8\'\'{urllib.parse.quote(filename)}'}
                )

    # Fallback: styled python-pptx generation
    prompt = PPT_PROMPT.format(
        context=context[:8000],
        n_slides=8,
        tone="default",
        verbosity="standard",
        language=cfg.language or "zh",
    )
    raw = await qwen_client.generate(
        [{"role": "system", "content": "Return only valid JSON."},
         {"role": "user", "content": prompt}]
    )

    raw = raw.strip()
    if raw.startswith("```"):
        raw = raw.split("\n", 1)[1].rsplit("```", 1)[0].strip()

    try:
        data = json_lib.loads(raw)
    except Exception:
        raise HTTPException(status_code=500, detail="Failed to parse slide structure from AI")

    if not isinstance(data, dict) or "slides" not in data or not isinstance(data.get("slides"), list):
        raise HTTPException(status_code=500, detail="AI returned unexpected slide structure")

    # Template color scheme (fallback uses general theme)
    theme = {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_bg": RGBColor(0x1A, 0x1A, 0x2E),
        "title_text": RGBColor(0xFF, 0xFF, 0xFF),
        "accent": RGBColor(0x00, 0x7A, 0xFF),
        "heading": RGBColor(0x1A, 0x1A, 0x2E),
        "body": RGBColor(0x33, 0x33, 0x33),
        "sub": RGBColor(0x66, 0x66, 0x66),
        "bar": RGBColor(0x00, 0x7A, 0xFF),
        "bullet_dot": RGBColor(0x00, 0x7A, 0xFF),
    }

    SLIDE_W = Inches(13.33)
    SLIDE_H = Inches(7.5)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    def _add_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_shape(slide, left, top, width, height, fill_color):
        from pptx.enum.shapes import MSO_SHAPE
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    def _add_text(slide, left, top, width, height, text, font_size, color, bold=False, alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
        return txBox

    # --- Title Slide ---
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(title_slide, theme["title_bg"])
    _add_shape(title_slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), theme["accent"])
    _add_text(
        title_slide, Inches(1.2), Inches(2.0), Inches(10.9), Inches(1.8),
        data.get("title", "Presentation"), 44, theme["title_text"], bold=True, alignment=PP_ALIGN.LEFT,
    )
    subtitle = data.get("subtitle", "Generated by Noteflow AI")
    _add_text(
        title_slide, Inches(1.2), Inches(3.8), Inches(10.9), Inches(0.8),
        subtitle, 20, theme["sub"], alignment=PP_ALIGN.LEFT,
    )
    _add_shape(title_slide, Inches(1.2), Inches(5.2), Inches(2.5), Inches(0.06), theme["accent"])
    _add_text(
        title_slide, Inches(1.2), Inches(6.2), Inches(4), Inches(0.5),
        "Generated by Noteflow AI", 12, theme["sub"],
    )

    # --- Content Slides ---
    for idx, slide_data in enumerate(data.get("slides", [])):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_bg(slide, theme["bg"])
        _add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), theme["bar"])
        _add_text(
            slide, Inches(11.5), Inches(0.3), Inches(1.5), Inches(0.4),
            f"{idx + 1:02d}", 14, theme["sub"], alignment=PP_ALIGN.RIGHT,
        )
        _add_shape(slide, Inches(0.8), Inches(0.9), Inches(0.06), Inches(0.7), theme["accent"])
        _add_text(
            slide, Inches(1.1), Inches(0.8), Inches(10), Inches(0.9),
            slide_data.get("title", ""), 32, theme["heading"], bold=True,
        )

        bullets = slide_data.get("bullets", [])
        y_offset = 2.0
        for bullet in bullets:
            if isinstance(bullet, dict):
                main_text = bullet.get("main", "")
                sub_text = bullet.get("sub", "")
            else:
                main_text = str(bullet)
                sub_text = ""
            _add_shape(slide, Inches(1.2), Inches(y_offset + 0.12), Inches(0.12), Inches(0.12), theme["bullet_dot"])
            _add_text(slide, Inches(1.6), Inches(y_offset - 0.05), Inches(10), Inches(0.5), main_text, 20, theme["body"], bold=True)
            y_offset += 0.45
            if sub_text:
                _add_text(slide, Inches(1.6), Inches(y_offset - 0.08), Inches(10), Inches(0.4), sub_text, 15, theme["sub"])
                y_offset += 0.38

        _add_shape(slide, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.01), theme["sub"])
        _add_text(slide, Inches(0.8), Inches(6.85), Inches(4), Inches(0.35), data.get("title", ""), 10, theme["sub"])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{urllib.parse.quote(filename)}"; filename*=UTF-8\'\'{urllib.parse.quote(filename)}'}
    )


@router.post("/podcast")
async def generate_podcast(
    notebook_id: uuid.UUID,
    db: AsyncSession = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    if not await permission_service.check_permission(db, notebook_id, current_user.id, "view"):
        raise HTTPException(status_code=403, detail="Forbidden")

    if not settings.ALIBABA_TTS_APPKEY or not settings.ALIBABA_TTS_TOKEN:
        raise HTTPException(status_code=501, detail="TTS service not configured")

    context = await _get_source_context(db, notebook_id)
    if not context:
        raise HTTPException(status_code=400, detail="No source documents available")

    # Generate dialogue
    prompt = PODCAST_PROMPT.format(context=context[:6000])
    dialogue = await qwen_client.generate(
        [{"role": "system", "content": "Generate engaging podcast dialogue. Follow the format exactly."},
         {"role": "user", "content": prompt}]
    )

    # Parse HOST:/GUEST: lines
    lines: list[tuple[str, str]] = []
    for line in dialogue.strip().split('\n'):
        line = line.strip()
        if line.upper().startswith('HOST:'):
            lines.append(('host', line[5:].strip()))
        elif line.upper().startswith('GUEST:'):
            lines.append(('guest', line[6:].strip()))

    if not lines:
        raise HTTPException(status_code=500, detail="Failed to generate dialogue structure")

    # Generate TTS for each line and merge
    host_voice = "ailun"    # Female Mandarin
    guest_voice = "aicheng"  # Male Mandarin
    silence = AudioSegment.silent(duration=400)  # 400ms pause between lines

    with tempfile.TemporaryDirectory() as tmpdir:
        segments: list[AudioSegment] = []
        for i, (speaker, text) in enumerate(lines):
            if not text:
                continue
            voice = host_voice if speaker == 'host' else guest_voice
            audio_bytes = await text_to_speech(
                text, voice,
                settings.ALIBABA_TTS_APPKEY,
                settings.ALIBABA_TTS_TOKEN,
            )
            seg_path = os.path.join(tmpdir, f"seg_{i}.mp3")
            with open(seg_path, 'wb') as f:
                f.write(audio_bytes)
            segments.append(AudioSegment.from_mp3(seg_path))
            segments.append(silence)

        if not segments:
            raise HTTPException(status_code=500, detail="No audio segments generated")

        combined = segments[0]
        for seg in segments[1:]:
            combined += seg

        out_path = os.path.join(tmpdir, "podcast.mp3")
        combined.export(out_path, format="mp3")
        with open(out_path, 'rb') as f:
            mp3_bytes = f.read()

    return StreamingResponse(
        io.BytesIO(mp3_bytes),
        media_type="audio/mpeg",
        headers={"Content-Disposition": 'attachment; filename="podcast.mp3"'}
    )


class StudioRequest(BaseModel):
    source_ids: list[str] | None = None


# ── Custom Skills CRUD (must be before /{content_type} catch-all) ──

from backend.models.custom_skill import CustomSkill
from backend.models.notebook_member import NotebookMember
from sqlalchemy import or_


class CustomSkillCreate(BaseModel):
    name: str
    prompt: str
    icon: str = "💡"
    all_notebooks: bool = True
    shared_with_team: bool = False


class CustomSkillUpdate(BaseModel):
    name: str | None = None
    prompt: str | None = None
    icon: str | None = None
    all_notebooks: bool | None = None
    shared_with_team: bool | None = None


class CustomSkillResponse(BaseModel):
    id: str
    name: str
    prompt: str
    icon: str
    created_by: str
    notebook_id: str
    all_notebooks: bool
    shared_with_team: bool

    model_config = {"from_attributes": True}


@router.get("/custom-skills")
async def list_custom_skills(
    notebook_id: str,
    user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """List custom skills visible in this notebook for the current user."""
    nb_uuid = uuid.UUID(notebook_id)

    # Find all members of the current notebook (including owner) to resolve team-shared skills
    member_stmt = select(NotebookMember.user_id).where(NotebookMember.notebook_id == nb_uuid)
    member_result = await db.execute(member_stmt)
    member_ids = {row[0] for row in member_result}
    # Also include the notebook owner (not always in notebook_members table)
    nb = await db.get(Notebook, nb_uuid)
    if nb:
        member_ids.add(nb.owner_id)

    # Skills visible to this user in this notebook:
    # 1. My skills with all_notebooks=True (visible everywhere)
    # 2. My skills created in this notebook
    # 3. Team-shared skills created in this notebook
    # 4. Team-shared skills where the creator is also a member of this notebook
    #    (covers the case where creator has all_notebooks=true and shared from another notebook)
    conditions = [
        (CustomSkill.created_by == user.id) & (CustomSkill.all_notebooks == True),
        (CustomSkill.created_by == user.id) & (CustomSkill.notebook_id == nb_uuid),
        (CustomSkill.notebook_id == nb_uuid) & (CustomSkill.shared_with_team == True),
    ]
    if member_ids:
        conditions.append(
            (CustomSkill.shared_with_team == True)
            & (CustomSkill.all_notebooks == True)
            & (CustomSkill.created_by.in_(member_ids))
        )

    stmt = select(CustomSkill).where(or_(*conditions))
    result = await db.execute(stmt)
    skills = result.scalars().all()

    seen = set()
    unique = []
    for s in skills:
        if s.id not in seen:
            seen.add(s.id)
            unique.append({
                "id": str(s.id), "name": s.name, "prompt": s.prompt, "icon": s.icon,
                "created_by": str(s.created_by), "notebook_id": str(s.notebook_id),
                "all_notebooks": s.all_notebooks, "shared_with_team": s.shared_with_team,
            })
    return unique


@router.post("/custom-skills")
async def create_custom_skill(
    notebook_id: str,
    body: CustomSkillCreate,
    user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Create a custom skill."""
    if not body.name.strip() or not body.prompt.strip():
        raise HTTPException(status_code=400, detail="Name and prompt are required")

    skill = CustomSkill(
        name=body.name.strip(),
        prompt=body.prompt.strip(),
        icon=body.icon or "💡",
        created_by=user.id,
        notebook_id=uuid.UUID(notebook_id),
        all_notebooks=body.all_notebooks,
        shared_with_team=body.shared_with_team,
    )
    db.add(skill)
    await db.commit()
    await db.refresh(skill)
    return {
        "id": str(skill.id), "name": skill.name, "prompt": skill.prompt, "icon": skill.icon,
        "created_by": str(skill.created_by), "notebook_id": str(skill.notebook_id),
        "all_notebooks": skill.all_notebooks, "shared_with_team": skill.shared_with_team,
    }


@router.patch("/custom-skills/{skill_id}")
async def update_custom_skill(
    notebook_id: str,
    skill_id: str,
    body: CustomSkillUpdate,
    user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Update a custom skill (creator or notebook owner/editor)."""
    skill = await db.get(CustomSkill, uuid.UUID(skill_id))
    if not skill:
        raise HTTPException(status_code=404, detail="Skill not found")
    # Allow creator, or notebook owner/editor for shared skills
    if skill.created_by != user.id:
        can_edit = await permission_service.check_permission(db, uuid.UUID(notebook_id), user.id, "rename")
        if not can_edit:
            raise HTTPException(status_code=403, detail="No permission to edit this skill")
    if body.name is not None:
        skill.name = body.name.strip()
    if body.prompt is not None:
        skill.prompt = body.prompt.strip()
    if body.icon is not None:
        skill.icon = body.icon
    if body.all_notebooks is not None:
        skill.all_notebooks = body.all_notebooks
    if body.shared_with_team is not None:
        skill.shared_with_team = body.shared_with_team
    await db.commit()
    return {"data": {"message": "Skill updated"}}


@router.delete("/custom-skills/{skill_id}")
async def delete_custom_skill(
    notebook_id: str,
    skill_id: str,
    user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Delete a custom skill (creator or notebook owner/editor)."""
    skill = await db.get(CustomSkill, uuid.UUID(skill_id))
    if not skill:
        raise HTTPException(status_code=404, detail="Skill not found")
    if skill.created_by != user.id:
        can_delete = await permission_service.check_permission(db, uuid.UUID(notebook_id), user.id, "rename")
        if not can_delete:
            raise HTTPException(status_code=403, detail="No permission to delete this skill")
    await db.delete(skill)
    await db.commit()
    return {"data": {"message": "Skill deleted"}}


@router.post("/custom-skills/{skill_id}/execute")
async def execute_custom_skill(
    notebook_id: str,
    skill_id: str,
    body: StudioRequest = Body(default=StudioRequest()),
    user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Execute a custom skill — same as generate_content but uses custom prompt."""
    skill = await db.get(CustomSkill, uuid.UUID(skill_id))
    if not skill:
        raise HTTPException(status_code=404, detail="Skill not found")

    source_ids = body.source_ids if body else None

    from backend.meeting.service import get_live_transcript_for_notebook_async
    live_transcript = await get_live_transcript_for_notebook_async(notebook_id)
    if live_transcript:
        context = live_transcript
    else:
        context = await _get_source_context(db, uuid.UUID(notebook_id), source_ids=source_ids)

    if not context:
        raise HTTPException(status_code=400, detail="No ready sources available")

    lang = _detect_language(context)
    prompt = skill.prompt + f"\n\nDOCUMENTS:\n{context}"
    system_msg = f"You are a helpful assistant. Write your entire response in {lang}."
    messages = [
        {"role": "system", "content": system_msg},
        {"role": "user", "content": prompt},
    ]
    content = await qwen_client.generate(messages)

    # Save as ChatMessage
    from backend.models.chat_message import ChatMessage
    lines = [l for l in content.split("\n") if l.strip() and not l.strip().startswith("#")]
    collapsed = lines[0][:150] if lines else content[:150]

    msg = ChatMessage(
        notebook_id=uuid.UUID(notebook_id),
        user_id=user.id,
        role="assistant",
        content=content,
        citations=[],
        msg_metadata={
            "type": "skill_output",
            "skill_type": str(skill.id),
            "skill_label": skill.name,
            "collapsed_summary": collapsed,
        },
    )
    db.add(msg)
    await db.commit()

    return {"content": content}


# ── Generic content generation (catch-all, must be LAST) ──────

@router.post("/{content_type}")
async def generate_content(
    notebook_id: str,
    content_type: str,
    body: StudioRequest = Body(default=StudioRequest()),
    user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
) -> dict:
    if not await permission_service.check_permission(db, uuid.UUID(notebook_id), user.id, "view"):
        raise HTTPException(status_code=403, detail="No access to this notebook")

    if content_type not in PROMPTS:
        raise HTTPException(status_code=400, detail=f"Invalid type. Allowed: {list(PROMPTS.keys())}")

    source_ids = body.source_ids if body else None
    logger.info("Studio %s: received source_ids=%s", content_type, source_ids)
    t_start = time.time()

    # If there's an active meeting, prioritize live transcript
    from backend.meeting.service import get_live_transcript_for_notebook_async
    live_transcript = await get_live_transcript_for_notebook_async(notebook_id)
    if live_transcript:
        context = live_transcript
        logger.info("Studio %s: using live meeting transcript (%d chars)", content_type, len(context))
    else:
        context = await _get_source_context(db, uuid.UUID(notebook_id), source_ids=source_ids)

    t_context = time.time()
    if not context:
        raise HTTPException(status_code=400, detail="No ready sources available for generation")

    logger.info("Studio %s: context retrieval %.1fs, %d chars", content_type, t_context - t_start, len(context))

    lang = _detect_language(context)
    prompt = PROMPTS[content_type].format(context=context)
    if content_type == "mindmap":
        system_msg = f"You are an expert at creating mind map structures from documents. Return ONLY valid JSON. You MUST write all labels in {lang}."
    else:
        system_msg = f"You are a helpful assistant that generates content from source documents. You MUST write your entire response in {lang}. Do NOT use any other language."
    messages = [
        {"role": "system", "content": system_msg},
        {"role": "user", "content": prompt},
    ]
    content = await qwen_client.generate(messages)
    t_llm = time.time()

    logger.info("Skill %s: LLM generation %.1fs, %d chars output. Total: %.1fs", content_type, t_llm - t_context, len(content), t_llm - t_start)

    # Save as ChatMessage for display in Chat panel
    from backend.models.chat_message import ChatMessage

    skill_labels = {
        "summary": "Summary", "faq": "FAQ", "action_items": "Action Items",
        "swot": "SWOT Analysis", "recommendations": "Recommendations",
        "risk_analysis": "Risk Analysis", "decision_support": "Decision Support",
        "study_guide": "Study Guide", "mindmap": "Mind Map",
    }
    if True:
        label = skill_labels.get(content_type, content_type.replace("_", " ").title())

        # Extract collapsed summary
        if content_type == "mindmap":
            collapsed = "Click to view mind map"
        else:
            lines = [l for l in content.split("\n") if l.strip() and not l.strip().startswith("#") and not l.strip().startswith("```")]
            collapsed = lines[0][:150] if lines else content[:100]

        msg = ChatMessage(
            notebook_id=uuid.UUID(notebook_id),
            user_id=user.id,
            role="assistant",
            content=content,
            citations=[],
            msg_metadata={
                "type": "skill_output",
                "skill_type": content_type,
                "skill_label": label,
                "collapsed_summary": collapsed,
            },
        )
        db.add(msg)
        await db.commit()

    return {"content": content}


